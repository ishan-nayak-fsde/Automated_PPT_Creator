import streamlit as st
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Function to call Ollama locally
def call_llm(prompt):
    response = requests.post(
        "http://localhost:11434/api/generate",
        json={
            "model": "llama3",
            "prompt": prompt,
            "stream": False
        }
    )
    return response.json()["response"]

# Function to generate PPT
def create_ppt(title, slide_content):
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title

    # Content Slides
    for content in slide_content:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = content["title"]
        slide.placeholders[1].text = content["content"]

    file_path = "generated_presentation.pptx"
    prs.save(file_path)
    return file_path

# Streamlit UI
st.title("🧠 Agentic AI PPT Generator")

subject = st.text_input("Enter Subject")
content = st.text_area("Enter Content")
use_ai = st.checkbox("Enhance with AI")

if st.button("Generate PPT"):
    if use_ai:
        prompt = f"""
        Create a structured PowerPoint outline.
        Subject: {subject}
        Content: {content}

        Format:
        Slide Title:
        Slide Content:
        """
        ai_output = call_llm(prompt)
    else:
        ai_output = f"Slide Title: {subject}\nSlide Content: {content}"

    slides = []
    parts = ai_output.split("Slide Title:")
    
    for part in parts[1:]:
        title, body = part.split("Slide Content:")
        slides.append({
            "title": title.strip(),
            "content": body.strip()
        })

    file_path = create_ppt(subject, slides)

    with open(file_path, "rb") as f:
        st.download_button(
            label="Download PPT",
            data=f,
            file_name="AI_Presentation.pptx"
        )
